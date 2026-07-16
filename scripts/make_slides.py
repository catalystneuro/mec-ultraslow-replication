"""Build a brief ~5-minute PPTX summarizing the replication findings.

Output uploads cleanly to Google Slides.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = f"{ROOT}/figures"
OUT = f"{ROOT}/MEC_ultraslow_replication.pptx"

# palette
INK = RGBColor(0x1A, 0x1A, 0x1A)
BLUE = RGBColor(0x21, 0x66, 0xAC)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
         space_after=8, bullet=False, level=0, italic=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.level = level
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Arial"
    return p


def bar(s, color, t=0.0, h=0.14):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(t),
                            SW, Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


# ---------------------------------------------------------------- 1. Title
s = slide()
bar(s, BLUE, 0, 0.22)
tf = box(s, 0.9, 2.2, 11.5, 3.0)
para(tf, "Do minute-scale oscillatory sequences in entorhinal cortex "
         "replicate in other datasets?", size=34, bold=True, color=INK, first=True)
para(tf, "A replication of Gonzalo Cogno et al. (2024), Nature 625:338–344, "
         "across public DANDI recordings", size=19, color=GREY, space_after=4)
tf2 = box(s, 0.9, 5.5, 11.5, 1.4)
para(tf2, "Prompted by dandi/helpdesk #156", size=15, color=GREY, first=True)
para(tf2, "CatalystNeuro", size=15, color=GREY)
bar(s, BLUE, 7.28, 0.14)

# ---------------------------------------------------------------- 2. Claim / question
s = slide()
bar(s, BLUE)
para(box(s, 0.7, 0.5, 12, 1), "The claim, and the question we can actually test",
     size=28, bold=True, first=True)
tf = box(s, 0.7, 1.7, 12.0, 5.4)
para(tf, "The finding", size=20, bold=True, color=BLUE, first=True, space_after=4)
para(tf, "In mouse MEC, activity organizes into ultraslow oscillations (periods "
         "of tens of seconds to minutes), and cells fire in periodic sequences "
         "during them.", size=18, bullet=True)
para(tf, "Observed in one condition only: head-fixed on a wheel, in darkness, no "
         "rewards. Darkness is where it was seen, not shown to be required.",
         size=18, bullet=True, space_after=14)
para(tf, "The authors' own open question", size=20, bold=True, color=BLUE, space_after=4)
para(tf, "\"It remains an open question whether the ultraslow oscillatory "
         "sequences are present across a broader spectrum of behaviours ... and "
         "in the presence of salient visual feedback.\"", size=17, italic=True,
         color=GREY, space_after=14)
para(tf, "So testing navigation datasets is a partial answer to their question, "
         "not a broken replication: do the sequences survive navigation and "
         "visual input, and do they appear outside the mouse?", size=18, bullet=True)

# ---------------------------------------------------------------- 3. Method + control
s = slide()
bar(s, BLUE)
para(box(s, 0.7, 0.5, 12, 1), "Method, and why a positive control was non-negotiable",
     size=28, bold=True, first=True)
tf = box(s, 0.7, 1.7, 5.9, 5.4)
para(tf, "The pipeline (per the paper)", size=20, bold=True, color=BLUE, first=True, space_after=4)
para(tf, "Bin spikes at 120 ms, smooth (σ = 5 s), binarize", size=17, bullet=True)
para(tf, "Detect ultraslow PSD peak below 0.1 Hz", size=17, bullet=True)
para(tf, "Population sequence test: PCA on time × cell, cells at staggered "
         "phases. Null = circular-shift each cell.", size=17, bullet=True)
para(tf, "Windowed (300 s), because sequences are intermittent", size=17,
     bullet=True, space_after=12)
para(tf, "Single-cell rhythmicity is NOT the test", size=18, bold=True, color=RED)
para(tf, "It is high everywhere (~50–75% of cells, even in V1). Only the "
         "population test discriminates, in the paper too.", size=16, color=GREY)

tf = box(s, 6.9, 1.7, 5.7, 5.4)
para(tf, "The keystone: a validated detector", size=20, bold=True, color=BLUE,
     first=True, space_after=4)
para(tf, "We obtained the paper's original Neuropixels wheel data (EBRAINS). A "
         "detector that cannot fire where the effect is known to exist cannot be "
         "trusted to report a null anywhere else.", size=17, bullet=True, space_after=10)
para(tf, "It fires on both original wheel mice (2/2). Two subtle statistical "
         "traps (band-power blindness, an FDR resolution floor) were caught only "
         "because the control existed.", size=17, bullet=True)

# ---------------------------------------------------------------- 4. Results (figure)
s = slide()
bar(s, RED)
para(box(s, 0.7, 0.42, 12, 0.9),
     "Result: sequences do not replicate in mouse MEC during navigation",
     size=26, bold=True, first=True)
s.shapes.add_picture(f"{FIG}/SUMMARY_all_sessions.png", Inches(0.35),
                     Inches(1.35), width=Inches(12.65))
tf = box(s, 0.7, 6.55, 12, 0.8)
para(tf, "Wheel/darkness 2/2  •  VR track 0/20  •  X-maze 3/114 (chance)  •  "
         "V1 region control 9/110, HIGHER than MEC  •  134 mouse MEC navigation "
         "sessions at chance", size=15, bold=True, color=INK, first=True)

# ---------------------------------------------------------------- 5. Confounds
s = slide()
bar(s, RED)
para(box(s, 0.7, 0.5, 12, 1),
     "The two apparent positives are behavioural confounds", size=26, bold=True, first=True)
para(box(s, 0.7, 1.4, 12, 0.6),
     "Each nailed by a dedicated control, in two species", size=17, color=GREY, first=True)
s.shapes.add_picture(f"{FIG}/macaque_task_control.png", Inches(6.85),
                     Inches(2.0), height=Inches(3.9))
tf = box(s, 0.7, 2.1, 5.9, 5.0)
para(tf, "Mouse visual cortex (9/110)", size=19, bold=True, color=INK, first=True, space_after=3)
para(tf, "The animal runs laps, the scene repeats, V1 tiles the lap cycle. "
         "Position, MEC-PC1 and V1-PC1 all peak at the same 0.007–0.010 Hz lap "
         "rhythm.", size=16, bullet=True, space_after=14)
para(tf, "Macaque EC (8/15)", size=19, bold=True, color=INK, space_after=3)
para(tf, "PC1 tracks trial-onset density: it drops to zero during task pauses "
         "and resumes with the trials. It is task-engagement block structure over "
         "a multi-hour session, not an intrinsic rhythm.", size=16, bullet=True)

# ---------------------------------------------------------------- 6. Takeaways
s = slide()
bar(s, BLUE)
para(box(s, 0.7, 0.5, 12, 1), "Takeaways", size=28, bold=True, first=True)
tf = box(s, 0.7, 1.7, 12.0, 5.4)
para(tf, "The pipeline reproduces the original effect on the wheel data, and "
         "sequences are not detectable in mouse MEC during navigation "
         "(0/20 + 3/114 sessions).", size=19, bullet=True, first=True, space_after=12)
para(tf, "This is consistent with the authors' own speculation that sequences "
         "may reset under strong landmarks or sensory input, a partial answer to "
         "their open question.", size=19, bullet=True, space_after=12)
para(tf, "It does NOT establish that darkness is required: the paper never ran "
         "another condition, and a null cannot separate \"sensory drive abolishes "
         "the effect\" from \"insufficient power\".", size=19, bullet=True, space_after=12)
para(tf, "Single-cell rhythmicity is non-specific: any replication that stops "
         "there reports a false positive. The population test is load-bearing.",
         size=19, bullet=True, space_after=12)
para(tf, "Depositing the EBRAINS wheel data on DANDI would be valuable: it is "
         "the only public recording in the condition where the effect is known "
         "to exist.", size=19, bullet=True)

prs.save(OUT)
print("wrote", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
